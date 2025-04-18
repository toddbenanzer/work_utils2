"""
Utilities for adding styled pandas DataFrames as tables to PowerPoint slides.
Supports row, column, and cell-specific styling, merged cells, and data formatting.
"""

import copy
from typing import Any, Dict, Optional

import numpy as np
import pandas as pd
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.slide import Slide
from pptx.table import Table, _Cell
from pptx.text.text import TextFrame

from .utils import (
    ALIGNMENT_MAP,
    VERTICAL_ALIGNMENT_MAP,
    Inches,
    Pt,
    deep_merge_dicts,
    get_color_from_value,
)


def get_default_table_style() -> Dict[str, Any]:
    """Returns a default style configuration dictionary for tables."""
    return {
        "position": {"left": 1.0, "top": 1.0},
        "dimensions": {"width": 0, "height": 0},
        "row_heights": 0.4,
        "col_widths": 1.2,
        "number_formats": {},
        "style_definitions": {
            "header": {
                "font_family": "Arial",
                "font_size": 12,
                "bold": True,
                "italic": False,
                "fill_color": (0, 112, 192),
                "font_color": (255, 255, 255),
                "h_align": "center",
                "v_align": "middle",
                "borders": {"bottom": {"width": 1.0, "color": (0, 0, 0)}},
            },
            "standard": {
                "font_family": "Arial",
                "font_size": 11,
                "bold": False,
                "italic": False,
                "fill_color": (255, 255, 255),
                "font_color": (0, 0, 0),
                "h_align": "left",
                "v_align": "middle",
                "borders": {"bottom": {"width": 1.0, "color": (0, 0, 0)}},
            },
            "subtotal": {
                "font_family": "Arial",
                "font_size": 11,
                "bold": True,
                "italic": False,
                "fill_color": (240, 240, 240),
                "font_color": (0, 0, 0),
                "h_align": "left",
                "v_align": "middle",
                "borders": {
                    "top": {"width": 1.0, "color": (0, 0, 0)},
                    "bottom": {"width": 1.0, "color": (0, 0, 0)},
                },
            },
            "total": {
                "font_family": "Arial",
                "font_size": 12,
                "bold": True,
                "italic": False,
                "fill_color": (217, 217, 217),
                "font_color": (0, 0, 0),
                "h_align": "left",
                "v_align": "middle",
                "borders": {
                    "top": {"width": 1.0, "color": (0, 0, 0)},
                    "bottom": {"width": 1.0, "color": (0, 0, 0)},
                },
            },
            "numeric_col": {"h_align": "right", "font_family": "Calibri"},
            "currency_col": {"h_align": "right", "font_family": "Calibri"},
            "category_col": {"bold": True, "h_align": "center"},
            "highlight_red": {"font_color": (200, 0, 0), "bold": True},
        },
        "row_styles": {0: "header"},
        "default_row_style": "standard",
        "column_styles": {},
        "cell_styles": {},
        "indent_cells": [],
        "indent_amount": 0.2,
        "merged_cells": [],
    }


def create_table_style(custom_style: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    """
    Creates a complete table style by merging custom settings with defaults.
    
    Parameters:
    -----------
    custom_style : dict, optional
        Custom style overrides to apply on top of the default style.
        
    Returns:
    --------
    dict
        A complete style settings dictionary.
    """
    base_style = get_default_table_style()
    if custom_style is None:
        return base_style
    return deep_merge_dicts(base_style, custom_style)


# --- Data Formatting Module ---
class DataFormatter:
    """Handles formatting of data values (initially numbers, extendable)."""

    @staticmethod
    def format_value(value: Any, format_spec: Dict[str, Any]) -> str:
        """
        Formats a single value according to the specified format.
        Currently handles numeric types, returns others as strings.
        
        Parameters:
        -----------
        value : Any
            The value to format.
        format_spec : dict
            Dictionary with formatting specifications (type, decimal_places, etc.).
            
        Returns:
        --------
        str
            The formatted string representation of the value.
        """
        if pd.isna(value) or value is None:
            return ""

        format_type = format_spec.get("type", "auto")

        # --- Numeric Formatting Logic ---
        is_numeric = False
        numeric_value = None
        try:
            # Check for numeric types, handling strings with commas
            if isinstance(value, (int, float, np.number)):
                numeric_value = float(value)
                is_numeric = True
            elif isinstance(value, str):
                cleaned_value = value.replace(",", "")
                if cleaned_value.replace(".", "", 1).replace("-", "", 1).isdigit():
                    numeric_value = float(cleaned_value)
                    is_numeric = True
        except (ValueError, TypeError):
            is_numeric = False  # Failed conversion

        if is_numeric and format_type in [
            "number",
            "currency",
            "percentage",
            "thousands",
            "millions",
        ]:
            decimal_places = format_spec.get("decimal_places", 2)
            show_commas = format_spec.get("show_commas", True)
            symbol = format_spec.get("symbol", "$")
            value_to_format = numeric_value
            suffix = ""
            prefix = ""

            if format_type == "percentage":
                value_to_format *= 100
                suffix = "%"
            elif format_type == "currency":
                prefix = symbol
            elif format_type == "thousands":
                value_to_format /= 1000
                suffix = "K"
            elif format_type == "millions":
                value_to_format /= 1000000
                suffix = "M"

            fmt_str = (
                f"{{:,.{decimal_places}f}}"
                if show_commas
                else f"{{:.{decimal_places}f}}"
            )
            formatted = fmt_str.format(value_to_format)

            # Remove trailing zeros and decimal point unless currency
            if "." in formatted and format_type != "currency":
                formatted = formatted.rstrip("0").rstrip(".")

            if format_type == "currency":
                # Handle negative currency format
                if numeric_value < 0:
                    return f"-{prefix}{formatted.lstrip('-')}"
                else:
                    return f"{prefix}{formatted}"
            return f"{formatted}{suffix}"

        # Default: return as string if no specific format applied
        return str(value)

    @staticmethod
    def apply_formats_to_dataframe(
        dataframe: pd.DataFrame, format_rules: Dict[str, Dict[str, Any]]
    ) -> pd.DataFrame:
        """
        Applies formatting to specified columns in a DataFrame.
        
        Parameters:
        -----------
        dataframe : pd.DataFrame
            The DataFrame to format.
        format_rules : dict
            Dictionary mapping column names to format specifications.
            
        Returns:
        --------
        pandas.DataFrame
            A new DataFrame with formatted values as strings.
        """
        if not format_rules:
            # Return copy with all columns as string if no rules
            return dataframe.astype(str)

        # Use deep copy to avoid modifying original df during processing
        df_formatted = dataframe.copy(deep=True)

        for col_name, format_spec in format_rules.items():
            if col_name in df_formatted.columns:
                df_formatted[col_name] = df_formatted[col_name].apply(
                    lambda x: DataFormatter.format_value(x, format_spec)
                )
            # No warning needed for missing columns

        # Ensure any columns *not* formatted are also strings for consistency
        for col in df_formatted.columns:
            if col not in format_rules:
                df_formatted[col] = df_formatted[col].astype(str)

        return df_formatted


# --- PowerPoint Table Formatter Class ---
class PowerPointTableFormatter:
    """
    Creates and styles a PowerPoint table from a pandas DataFrame.
    Supports row, column, and cell-specific styling, merged cells, data formatting.
    """

    def __init__(self, slide: Slide, dataframe: pd.DataFrame, style: Dict[str, Any]):
        """
        Initializes the formatter.
        
        Parameters:
        -----------
        slide : pptx.slide.Slide
            The PowerPoint slide to add the table to.
        dataframe : pd.DataFrame
            The data to populate the table. Must contain data convertible to string.
            Headers are derived from dataframe columns. Data starts from row 1.
        style : dict
            A style configuration dictionary (use create_table_style()).
        """
        if not isinstance(dataframe, pd.DataFrame):
            raise TypeError("Input 'dataframe' must be a pandas DataFrame.")
        if not isinstance(style, dict):
            raise TypeError("Input 'style' must be a dictionary.")

        self.slide = slide
        self.original_dataframe = dataframe.copy()
        self.style = style
        self.n_rows = dataframe.shape[0] + 1  # +1 for header row
        self.n_cols = dataframe.shape[1]
        self.table: Optional[Table] = None

        # Extract frequently used style parts for cleaner access
        self.style_definitions = self.style.get("style_definitions", {})
        self.row_style_map = self.style.get("row_styles", {0: "header"})
        self.default_row_style_key = self.style.get("default_row_style", "standard")
        self.column_style_map = self.style.get("column_styles", {})
        self.cell_style_map = self.style.get("cell_styles", {})
        self.merged_cells_list = self.style.get("merged_cells", [])
        self.indent_cells_list = self.style.get("indent_cells", [])
        self.indent_amount_val = self.style.get("indent_amount", 0.2)

        # Prepare data immediately - converts df to strings with formatting
        self.formatted_dataframe = self._prepare_data()

    def _prepare_data(self) -> pd.DataFrame:
        """Applies data formatting based on style configuration."""
        format_rules = self.style.get("number_formats", {})
        return DataFormatter.apply_formats_to_dataframe(
            self.original_dataframe, format_rules
        )

    def _create_table(self) -> Table:
        """Creates the table shape on the slide with calculated dimensions."""
        pos = self.style.get("position", {"left": 1.0, "top": 1.0})
        dim_override = self.style.get("dimensions", {"width": 0, "height": 0})
        col_widths = self.style.get("col_widths", 1.2)
        row_heights = self.style.get("row_heights", 0.4)

        # Determine Width
        if dim_override.get("width", 0) > 0:
            final_width = Inches(dim_override["width"])
        elif isinstance(col_widths, list):
            final_width = Inches(sum(w for w in col_widths if w > 0))
        elif col_widths > 0:
            final_width = Inches(col_widths * self.n_cols)
        else:
            final_width = Inches(1.2 * self.n_cols)  # Default fallback width

        # Determine Height
        if dim_override.get("height", 0) > 0:
            final_height = Inches(dim_override["height"])
        elif isinstance(row_heights, list):
            final_height = Inches(sum(h for h in row_heights if h > 0))
        elif row_heights > 0:
            final_height = Inches(row_heights * self.n_rows)
        else:
            final_height = Inches(0.4 * self.n_rows)  # Default fallback height

        table_shape = self.slide.shapes.add_table(
            self.n_rows,
            self.n_cols,
            Inches(pos.get("left", 1.0)),
            Inches(pos.get("top", 1.0)),
            final_width,
            final_height,
        )
        return table_shape.table

    def _apply_merges(self):
        """Applies cell merges as defined in the style configuration."""
        if self.table is None or not self.merged_cells_list:
            return

        processed_cells = set()  # Track top-left cells involved in a merge

        for r1, c1, r2, c2 in self.merged_cells_list:
            # Basic coordinate validation
            if not (
                0 <= r1 < self.n_rows
                and 0 <= c1 < self.n_cols
                and r1 <= r2 < self.n_rows
                and c1 <= c2 < self.n_cols
            ):
                continue

            # Check if start cell is already part of another merge initiated earlier
            if (r1, c1) in processed_cells:
                continue

            cell1 = self.table.cell(r1, c1)
            # Check if the target cell is already spanned
            if cell1.is_spanned:
                continue

            # Proceed with merge if start cell is not spanned
            cell2 = self.table.cell(r2, c2)
            cell1.merge(cell2)
            processed_cells.add((r1, c1))  # Mark this start cell as processed

    def _populate_table(self):
        """Fills the table with headers and formatted data, respecting merged cells."""
        if self.table is None:
            raise ValueError("Table object is None, cannot populate.")

        # Add column headers (using original dataframe for headers)
        for c_idx, col_name in enumerate(self.original_dataframe.columns):
            cell = self.table.cell(0, c_idx)
            # Only write text to the 'master' cell (top-left) of a merged range
            if not cell.is_spanned:
                cell.text = str(col_name)

        # Add formatted data rows (using formatted_dataframe)
        for r_idx_data, row_data in enumerate(
            self.formatted_dataframe.itertuples(index=False)
        ):
            r_idx_table = r_idx_data + 1  # Table row index (+1 for header)
            for c_idx, value in enumerate(row_data):
                cell = self.table.cell(r_idx_table, c_idx)
                # Only write text to the 'master' cell (top-left) of a merged range
                if not cell.is_spanned:
                    cell.text = str(value)  # Ensure string conversion

    def _apply_dimensions(self):
        """Applies row heights and column widths if specified in the style."""
        if self.table is None:
            raise ValueError("Table object is None, cannot apply dimensions.")

        row_heights = self.style.get("row_heights", None)
        col_widths = self.style.get("col_widths", None)

        # Apply row heights
        if row_heights is not None:
            height_list = (
                [row_heights] * self.n_rows
                if isinstance(row_heights, (int, float))
                else row_heights
            )
            if len(height_list) > 0:  # Ensure list is not empty
                for i in range(self.n_rows):
                    try:
                        height_val = height_list[i % len(height_list)]  # Cycle if needed
                        if height_val > 0:
                            self.table.rows[i].height = Inches(height_val)
                    except IndexError:
                        pass  # Should not happen with modulo

        # Apply column widths
        if col_widths is not None:
            width_list = (
                [col_widths] * self.n_cols
                if isinstance(col_widths, (int, float))
                else col_widths
            )
            if len(width_list) > 0:
                for i in range(self.n_cols):
                    try:
                        width_val = width_list[i % len(width_list)]  # Cycle if needed
                        if width_val > 0:
                            self.table.columns[i].width = Inches(width_val)
                    except IndexError:
                        pass

    def _apply_font_style(self, text_frame: TextFrame, style: Dict[str, Any]):
        """Applies font settings from style dict to all runs in all paragraphs."""
        font_settings = {
            "name": style.get("font_family"),
            "size": Pt(style["font_size"]) if "font_size" in style else None,
            "bold": style.get("bold"),
            "italic": style.get("italic"),
            "color_rgb": (
                RGBColor(*style["font_color"]) if "font_color" in style else None
            ),
        }
        # Filter out None values
        active_font_settings = {k: v for k, v in font_settings.items() if v is not None}

        if not active_font_settings:
            return  # Nothing to apply

        for paragraph in text_frame.paragraphs:
            # Ensure at least one run exists to apply font to
            if not paragraph.runs:
                paragraph.add_run()

            for run in paragraph.runs:
                font = run.font
                if "name" in active_font_settings:
                    font.name = active_font_settings["name"]
                if "size" in active_font_settings:
                    font.size = active_font_settings["size"]
                if "bold" in active_font_settings:
                    font.bold = active_font_settings["bold"]
                if "italic" in active_font_settings:
                    font.italic = active_font_settings["italic"]
                if "color_rgb" in active_font_settings:
                    font.color.rgb = active_font_settings["color_rgb"]

    def _apply_cell_style(self, cell: _Cell, cell_style: Dict[str, Any]):
        """Applies comprehensive styling (fill, align, font, borders) to a single cell."""
        # Do not style spanned cells directly; their appearance is controlled by the master cell.
        if cell.is_spanned:
            return

        # Cell Fill
        if "fill_color" in cell_style:
            r, g, b = cell_style["fill_color"]
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(r, g, b)
        else:
            cell.fill.background()  # Explicitly set background fill (no color)

        # Vertical Alignment
        if "v_align" in cell_style and cell_style["v_align"] in VERTICAL_ALIGNMENT_MAP:
            cell.vertical_anchor = VERTICAL_ALIGNMENT_MAP[cell_style["v_align"]]

        # Text Frame / Paragraph / Font Formatting
        text_frame = cell.text_frame
        text_frame.word_wrap = cell_style.get("word_wrap", True)

        # Apply paragraph settings (like alignment)
        if not text_frame.paragraphs:
            text_frame.add_paragraph()  # Ensure paragraph exists
        for paragraph in text_frame.paragraphs:
            if "h_align" in cell_style and cell_style["h_align"] in ALIGNMENT_MAP:
                paragraph.alignment = ALIGNMENT_MAP[cell_style["h_align"]]

        # Apply font settings
        self._apply_font_style(text_frame, cell_style)

        # Border Formatting
        self._apply_border_formatting(cell, cell_style.get("borders", {}))

    def _apply_border_formatting(
        self, cell: _Cell, border_style: Dict[str, Dict[str, Any]]
    ):
        """
        Applies border styles defined in border_style dict to a cell
        using direct OXML manipulation based on common patterns.
        """
        # Skip spanned cells; borders controlled by the master cell.
        if cell.is_spanned:
            return

        # Helper: Apply properties to an OXML line element (e.g., <a:lnT>)
        def apply_line_properties(line_element, side_style: Dict[str, Any]):
            """Sets properties on an OXML line element (e.g., a:lnT)."""

            # Helper to remove existing fill/dash child elements
            def clear_line_fill_and_dash(ln_el):
                # List of potential child tags related to fill and dash style
                tags_to_remove = [
                    "a:noFill",
                    "a:solidFill",
                    "a:gradFill",
                    "a:pattFill",  # Fill types
                    "a:prstDash",
                    "a:custDash",  # Dash types
                ]
                for tag in tags_to_remove:
                    # Find the element using the qualified tag name
                    child = ln_el.find(qn(tag))
                    # If found, remove it from its parent (ln_el)
                    if child is not None:
                        ln_el.remove(child)

            if not side_style:
                # If no style defined for this side, ensure no visible line
                clear_line_fill_and_dash(line_element)  # Remove fill/dash elements
                line_element.set("w", "0")  # Set width attribute to 0
                if "cap" in line_element.attrib:  # Remove cap attribute if present
                    del line_element.attrib["cap"]
                return

            # Width
            width_pts = side_style.get("width", 0)  # Default to 0 if not specified
            if not isinstance(width_pts, (int, float)):
                width_pts = 0

            if width_pts > 0:
                line_element.set("w", str(Pt(width_pts).emu))  # Width attribute in EMU
                line_element.set("cap", "sq")  # Line cap: 'sq' for square (solid look)
            else:
                # Width is 0 or invalid, ensure no visible line
                clear_line_fill_and_dash(line_element)  # Remove fill/dash
                line_element.set("w", "0")
                if "cap" in line_element.attrib:  # Remove cap
                    del line_element.attrib["cap"]
                return  # No need to process color/dash if width is 0

            # Clear existing fill/dash before applying potentially new ones
            clear_line_fill_and_dash(line_element)

            # Color (Only apply if width > 0)
            color_rgb = side_style.get("color")
            if (
                color_rgb is not None
                and isinstance(color_rgb, tuple)
                and len(color_rgb) == 3
            ):
                # Ensure values are valid RGB 0-255
                r, g, b = [max(0, min(255, c)) for c in color_rgb]

                # Create <a:solidFill><a:srgbClr val="RRGGBB"/></a:solidFill>
                solidFill = OxmlElement("a:solidFill")
                srgbClr = OxmlElement("a:srgbClr")

                # Format RGB tuple into "RRGGBB" hex string
                hex_color_string = f"{r:02x}{g:02x}{b:02x}"
                srgbClr.set("val", hex_color_string.upper())

                solidFill.append(srgbClr)
                line_element.append(solidFill)

            # Dash Style (Only apply if width > 0)
            dash_style_val = side_style.get("dash_style")
            if isinstance(dash_style_val, str) and dash_style_val:
                # Add the new dash style
                prstDash = OxmlElement("a:prstDash")
                prstDash.set("val", dash_style_val)
                line_element.append(prstDash)

        # Get or Add Generic OXML Element Helper
        def get_or_add(parent_element, tag):
            """Gets the child element with `tag` or creates, appends, and returns it."""
            element = parent_element.find(qn(tag))
            if element is None:
                element = OxmlElement(tag)
                parent_element.append(element)
            return element

        # Main Logic
        tcPr = cell._tc.get_or_add_tcPr()  # Get <a:tcPr> element

        # Define border tags
        border_tags = {
            "top": "a:lnT",
            "right": "a:lnR",
            "bottom": "a:lnB",
            "left": "a:lnL",
        }

        # Apply style for each specified border side
        for side, tag in border_tags.items():
            if side in border_style:
                line_element = get_or_add(tcPr, tag)
                apply_line_properties(line_element, border_style[side])

    def _apply_table_styling(self):
        """
        Applies styling using a hierarchical approach:
        1. Base Row Style
        2. Column Style Override (merged onto row style)
        3. Cell Specific Override (merged onto the result)
        """
        if self.table is None:
            raise ValueError("Table object is None, cannot apply styles.")

        # Get default row style definition once, ensure it's a dict
        default_row_style = self.style_definitions.get(self.default_row_style_key, {})
        if not isinstance(default_row_style, dict):
            default_row_style = {}

        # Cache looked-up style definitions to avoid repeated lookups
        resolved_style_defs = {self.default_row_style_key: default_row_style}

        for r_idx in range(self.n_rows):
            # 1. Determine Base Row Style
            row_style_key = self.row_style_map.get(r_idx, self.default_row_style_key)
            if row_style_key not in resolved_style_defs:
                resolved_style_defs[row_style_key] = self.style_definitions.get(
                    row_style_key, {}
                )
                if not resolved_style_defs[row_style_key]:
                    resolved_style_defs[row_style_key] = default_row_style  # Fallback to default

            base_row_style = resolved_style_defs[row_style_key]

            for c_idx in range(self.n_cols):
                cell = self.table.cell(r_idx, c_idx)
                # Skip styling spanned cells directly
                if cell.is_spanned:
                    continue

                # Start with a deep copy of the base row style for this cell
                current_cell_style = copy.deepcopy(base_row_style)

                # 2. Apply Column Style Override
                col_style_key = self.column_style_map.get(c_idx)
                if col_style_key:
                    if col_style_key not in resolved_style_defs:
                        resolved_style_defs[col_style_key] = self.style_definitions.get(
                            col_style_key, {}
                        )

                    col_style_def = resolved_style_defs[col_style_key]
                    if col_style_def:  # Apply only if found
                        current_cell_style = deep_merge_dicts(
                            current_cell_style, col_style_def
                        )

                # 3. Apply Cell Specific Override
                cell_override_style = self.cell_style_map.get((r_idx, c_idx))
                if cell_override_style and isinstance(cell_override_style, dict):
                    current_cell_style = deep_merge_dicts(
                        current_cell_style, cell_override_style
                    )

                # 4. Apply the final merged style to the cell
                self._apply_cell_style(cell, current_cell_style)

    def _apply_indentation(self):
        """Applies text indentation to specified cells, respecting merged cells."""
        if self.table is None or not self.indent_cells_list or self.indent_amount_val <= 0:
            return

        indent_val = Inches(self.indent_amount_val)

        for row_idx, col_idx in self.indent_cells_list:
            # Validate indices
            if not (0 <= row_idx < self.n_rows and 0 <= col_idx < self.n_cols):
                continue

            cell = self.table.cell(row_idx, col_idx)
            # Only apply indentation to the 'master' cell (top-left) of a merged range
            if not cell.is_spanned:
                text_frame = cell.text_frame
                if not text_frame.paragraphs:
                    text_frame.add_paragraph()  # Ensure paragraph exists
                for paragraph in text_frame.paragraphs:
                    # Set indent on paragraph format
                    pf = paragraph.paragraph_format
                    # Ensure level 0 for standard indent
                    if paragraph.level != 0:
                        paragraph.level = 0
                    pf.left_indent = indent_val

    def add_to_slide(self) -> Table:
        """
        Executes the full process to create, populate, merge, and style the table.
        
        Returns:
        --------
        pptx.table.Table
            The created and formatted PowerPoint table object.
        """
        self.table = self._create_table()
        self._apply_merges()  # Merge cells first
        self._populate_table()  # Then populate data
        self._apply_dimensions()  # Then set dimensions
        self._apply_table_styling()  # Then apply styles
        self._apply_indentation()  # Finally apply indentation
        return self.table


def add_table(
    slide: Slide, dataframe: pd.DataFrame, style_config: Optional[Dict[str, Any]] = None
) -> Table:
    """
    High-level function to add a pandas DataFrame as a styled table to a PowerPoint slide.
    
    Parameters:
    -----------
    slide : pptx.slide.Slide
        The PowerPoint slide object.
    dataframe : pd.DataFrame
        The data to display. Column names become headers.
    style_config : dict, optional
        A dictionary with custom style settings. If None, default styles are used.
        
    Returns:
    --------
    pptx.table.Table
        The created PowerPoint table object.
    """
    # Create the final style dictionary
    final_style = create_table_style(custom_style=style_config)
    
    # Instantiate the formatter and run the process
    formatter = PowerPointTableFormatter(slide, dataframe, final_style)
    return formatter.add_to_slide()
