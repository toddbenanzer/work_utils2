"""
Utilities for adding styled text boxes to PowerPoint slides.

Provides functions to create and style text boxes based on a list of strings,
using a configurable dictionary for styling options.
"""

from typing import Any, Dict, List, Optional, Union

from pptx.dml.color import RGBColor
from pptx.shapes.autoshape import Shape
from pptx.slide import Slide
from pptx.text.text import Font

from .utils import (
    ALIGNMENT_MAP,
    AUTOSIZE_MAP,
    VERTICAL_ALIGNMENT_MAP,
    Inches,
    Pt,
    apply_font_settings,
    deep_merge_dicts,
    get_color_from_value,
)


def get_default_textbox_style() -> Dict[str, Any]:
    """Returns a default style configuration dictionary for text boxes."""
    return {
        "position": {"left": 1.0, "top": 1.0},
        "dimensions": {"width": 6.0, "height": 1.0},
        "font": {
            "name": None,  # Use theme default
            "size": 12,
            "bold": False,
            "italic": False,
            "color": "000000",  # Default black as hex
        },
        "paragraph": {
            "alignment": "left",
            "space_before": 0,
            "space_after": 0,
            "line_spacing": 1.0,
        },
        "bullets": {
            "enabled": False,
            "level": 1,
        },
        "fill": {
            "type": "none",
        },
        "border": {
            "visible": False,
        },
        "auto_size": "shape_to_fit_text",
        "margins": {"left": 0.1, "right": 0.1, "top": 0.05, "bottom": 0.05},
        "vertical_anchor": "top",
    }


def create_textbox_style(
    custom_style: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    """
    Creates a complete textbox style by merging custom settings with defaults.

    Parameters:
    -----------
    custom_style : dict, optional
        Custom style overrides to apply on top of the default style.

    Returns:
    --------
    dict
        A complete style settings dictionary.
    """
    base_style = get_default_textbox_style()
    if custom_style is None:
        return base_style
    if not isinstance(custom_style, dict):
        raise TypeError("custom_style must be a dictionary or None.")
    return deep_merge_dicts(base_style, custom_style)


class PowerPointTextboxFormatter:
    """
    Creates and styles a PowerPoint textbox from a list of strings based on a style config.
    """

    def __init__(self, slide: Slide, text_content: List[str], style: Dict[str, Any]):
        """
        Initializes the textbox formatter.

        Parameters:
        -----------
        slide : pptx.slide.Slide
            The PowerPoint slide to add the textbox to.
        text_content : List[str]
            A list of strings, where each string becomes a paragraph.
        style : dict
            A complete style configuration dictionary (result of create_textbox_style()).
        """
        if not isinstance(text_content, list):
            raise TypeError("Input 'text_content' must be a list of strings.")
        if not isinstance(style, dict):
            raise TypeError("Input 'style' must be a dictionary.")

        self.slide = slide
        self.text_content = [str(item) for item in text_content]
        self.style = style
        self._textbox_shape: Optional[Shape] = None
        self.text_frame = None

    def _create_textbox_object(self) -> None:
        """Creates the textbox shape on the slide."""
        pos_cfg = self.style.get("position", {})
        dim_cfg = self.style.get("dimensions", {})

        left = Inches(float(pos_cfg.get("left", 1.0)))
        top = Inches(float(pos_cfg.get("top", 1.0)))
        width = Inches(float(dim_cfg.get("width", 6.0)))
        height = Inches(float(dim_cfg.get("height", 1.0)))

        self._textbox_shape = self.slide.shapes.add_textbox(left, top, width, height)
        self.text_frame = self._textbox_shape.text_frame

    def _apply_textbox_properties(self) -> None:
        """Applies shape-level properties like fill, border, margins, auto-size."""
        if not self._textbox_shape or not self.text_frame:
            raise RuntimeError("Textbox shape or text frame not initialized.")

        # Auto Size
        auto_size_cfg = self.style.get("auto_size", "none")
        auto_size_enum = AUTOSIZE_MAP.get(str(auto_size_cfg).lower())
        if auto_size_enum is None:
            raise ValueError(
                f"Invalid auto_size value: '{auto_size_cfg}'. Use one of {list(AUTOSIZE_MAP.keys())}"
            )
        self.text_frame.auto_size = auto_size_enum

        # Margins
        margin_cfg = self.style.get("margins", {})
        if "left" in margin_cfg:
            self.text_frame.margin_left = Inches(float(margin_cfg["left"]))
        if "right" in margin_cfg:
            self.text_frame.margin_right = Inches(float(margin_cfg["right"]))
        if "top" in margin_cfg:
            self.text_frame.margin_top = Inches(float(margin_cfg["top"]))
        if "bottom" in margin_cfg:
            self.text_frame.margin_bottom = Inches(float(margin_cfg["bottom"]))

        # Fill
        fill_cfg = self.style.get("fill", {})
        fill_type = str(fill_cfg.get("type", "none")).lower()
        shape_fill = self._textbox_shape.fill
        if fill_type == "none":
            shape_fill.background()
        elif fill_type == "solid":
            fill_color = fill_cfg.get("color")
            if fill_color is None:
                raise ValueError("Solid fill specified but no 'color' provided in fill style.")
            shape_fill.solid()
            shape_fill.fore_color.rgb = get_color_from_value(fill_color)
        else:
            raise ValueError(f"Unsupported fill type: '{fill_type}'. Use 'none' or 'solid'.")

        # Border (Line)
        border_cfg = self.style.get("border", {})
        is_visible = bool(border_cfg.get("visible", False))
        shape_line = self._textbox_shape.line
        if is_visible:
            border_color = border_cfg.get("color")
            border_width = border_cfg.get("width", 1.0)
            if border_color is None:
                raise ValueError("Border visibility is true but no 'color' provided in border style.")
            shape_line.color.rgb = get_color_from_value(border_color)
            shape_line.width = Pt(float(border_width))
        else:
            shape_line.fill.background()  # Makes line invisible

        # Vertical Alignment
        valign_cfg = self.style.get("vertical_anchor", "top").lower()
        valign_enum = VERTICAL_ALIGNMENT_MAP.get(valign_cfg)
        if valign_enum is None:
            raise ValueError(
                f"Invalid vertical_anchor value: '{valign_cfg}'. Use 'top', 'middle', or 'bottom'."
            )
        self.text_frame.vertical_anchor = valign_enum

    def _apply_text_and_paragraphs(self) -> None:
        """Adds text content and applies paragraph/font styles."""
        if not self.text_frame:
            raise RuntimeError("Text frame not initialized.")

        # Clear any default paragraph/text in the newly added textbox
        self.text_frame.clear()

        font_style = self.style.get("font", {})
        para_style = self.style.get("paragraph", {})
        bullet_style = self.style.get("bullets", {})
        bullets_enabled = bool(bullet_style.get("enabled", False))

        # Get Paragraph Settings
        align_str = str(para_style.get("alignment", "left")).lower()
        alignment = ALIGNMENT_MAP.get(align_str)
        if alignment is None:
            raise ValueError(
                f"Invalid paragraph alignment: '{align_str}'. Use one of {list(ALIGNMENT_MAP.keys())}"
            )

        space_before = Pt(int(para_style.get("space_before", 0)))
        space_after = Pt(int(para_style.get("space_after", 0)))
        line_spacing = float(para_style.get("line_spacing", 1.0))

        # Get Bullet Settings (if enabled)
        bullet_level = 0  # Default to 0 (no indent/bullet)
        if bullets_enabled:
            # Config uses 1-based level, pptx uses 0-based
            level_1_based = int(bullet_style.get("level", 1))
            if level_1_based < 1 or level_1_based > 9:  # pptx supports levels 0-8
                raise ValueError("Bullet level must be between 1 and 9.")
            bullet_level = level_1_based - 1

        # Add Paragraphs and Apply Styles
        for line_text in self.text_content:
            p = self.text_frame.add_paragraph()
            p.text = line_text

            # Apply paragraph settings
            p.alignment = alignment
            p.space_before = space_before
            p.space_after = space_after
            p.line_spacing = line_spacing

            # Apply font settings
            apply_font_settings(p.font, font_style)

            # Apply bullets if enabled
            if bullets_enabled:
                p.level = bullet_level

    def add_to_slide(self) -> Shape:
        """
        Executes the full process: creates the textbox shape and applies all styles.

        Returns:
        --------
        pptx.shapes.autoshape.Shape
            The created and styled PowerPoint textbox shape object.
        """
        self._create_textbox_object()
        self._apply_textbox_properties()
        self._apply_text_and_paragraphs()

        if not self._textbox_shape:
            raise RuntimeError("Textbox shape was not created successfully.")

        return self._textbox_shape


def add_textbox(
    slide: Slide,
    text_content: Union[str, List[str]],
    style_config: Optional[Dict[str, Any]] = None,
) -> Shape:
    """
    High-level function to add a styled textbox to a PowerPoint slide.

    Parameters:
    -----------
    slide : pptx.slide.Slide
        The PowerPoint slide object.
    text_content : Union[str, List[str]]
        The text to display. If a single string, it's treated as one paragraph.
        If a list of strings, each string becomes a separate paragraph.
    style_config : dict, optional
        A dictionary with custom style settings. These will be merged with
        the default style settings.

    Returns:
    --------
    pptx.shapes.autoshape.Shape
        The created and styled PowerPoint textbox shape object.
    """
    # Ensure text_content is a list
    if isinstance(text_content, str):
        content_list = [text_content]
    elif isinstance(text_content, list):
        content_list = [str(item) for item in text_content]
    else:
        raise TypeError("text_content must be a string or a list of strings.")

    # Create the final style dictionary
    final_style = create_textbox_style(custom_style=style_config)

    # Instantiate the formatter
    formatter = PowerPointTextboxFormatter(slide, content_list, final_style)

    # Run the creation and styling process
    return formatter.add_to_slide()
