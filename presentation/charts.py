"""
Utilities for adding pandas DataFrame charts to PowerPoint slides.

Provides functions to create and style charts based on pandas DataFrames,
using a configurable dictionary for styling options.
"""

from typing import Any, Dict, List, Optional, Union

import pandas as pd
from pptx.chart.chart import Chart
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.shapes.graphfrm import GraphicFrame
from pptx.slide import Slide

from .utils import (
    CHART_TYPE_MAP,
    DEFAULT_CHART_COLORS,
    LABEL_POSITION_MAP,
    LEGEND_POSITION_MAP,
    TICK_MARK_MAP,
    Inches,
    Pt,
    apply_font_settings,
    deep_merge_dicts,
    get_color_from_value,
)


def get_default_chart_style() -> Dict[str, Any]:
    """Returns a simplified default style configuration dictionary for charts."""
    return {
        "chart_type": "stacked_bar",
        "position": {"left": 1.0, "top": 1.5},
        "dimensions": {"width": 8.0, "height": 5.0},
        "title": {
            "text": None,
            "visible": True,
            "font": {"size": 14, "bold": True},
        },
        "data_labels": {
            "enabled": True,
            "position": "center",
            "font": {
                "size": 9,
                "bold": False,
                "color": (255, 255, 255),
            },
        },
        "legend": {
            "enabled": True,
            "position": "bottom",
            "font": {"size": 10},
        },
        "category_axis": {
            "visible": True,
            "tick_labels": {"visible": True, "font": {"size": 10}},
        },
        "value_axis": {
            "visible": True,
            "major_gridlines": True,
            "tick_labels": {"visible": True, "font": {"size": 10}},
        },
        "gap_width": 150,
        "overlap": 100,
        "colors": {},
    }


def create_chart_style(custom_style: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    """
    Creates a complete chart style by merging custom settings with defaults.
    
    Parameters:
    -----------
    custom_style : dict, optional
        Custom style overrides to apply on top of the default style.
        
    Returns:
    --------
    dict
        A complete style settings dictionary.
    """
    base_style = get_default_chart_style()
    if custom_style is None:
        return base_style
    if not isinstance(custom_style, dict):
        raise TypeError("custom_style must be a dictionary or None.")
    return deep_merge_dicts(base_style, custom_style)


class PowerPointChartFormatter:
    """
    Creates and styles a PowerPoint chart from a pandas DataFrame based on a style config.
    """

    def __init__(self, slide: Slide, dataframe: pd.DataFrame, style: Dict[str, Any]):
        """
        Initializes the chart formatter.
        
        Parameters:
        -----------
        slide : pptx.slide.Slide
            The PowerPoint slide to add the chart to.
        dataframe : pd.DataFrame
            The data to visualize. Assumes the first column contains categories
            and subsequent columns contain series data.
        style : dict
            A complete style configuration dictionary (result of create_chart_style()).
        """
        if not isinstance(dataframe, pd.DataFrame):
            raise TypeError("Input 'dataframe' must be a pandas DataFrame.")
        if dataframe.empty:
            raise ValueError("Input 'dataframe' cannot be empty.")
        if dataframe.shape[1] < 2:
            raise ValueError("DataFrame must have at least two columns (categories + 1 series).")
        if not isinstance(style, dict):
            raise TypeError("Input 'style' must be a dictionary.")

        self.slide = slide
        self.dataframe = dataframe.copy()
        self.style = style
        self.chart: Optional[Chart] = None
        self._chart_shape: Optional[GraphicFrame] = None

    def _prepare_chart_data(self) -> CategoryChartData:
        """Prepares the CategoryChartData object from the DataFrame."""
        chart_data = CategoryChartData()

        # Assume first column is category labels
        categories_col = self.dataframe.columns[0]
        # Ensure categories are strings for pptx
        chart_data.categories = self.dataframe[categories_col].astype(str).tolist()

        # Subsequent columns are series
        data_cols = self.dataframe.columns[1:]
        for col in data_cols:
            # Convert to numeric, coercing errors. Fill resulting NaN with None.
            series_data = pd.to_numeric(self.dataframe[col], errors="coerce")
            # Convert pandas NA (pd.NA or np.nan) to None for pptx
            series_data_list: List[Union[float, int, None]] = [
                val if pd.notna(val) else None for val in series_data
            ]
            chart_data.add_series(str(col), series_data_list)

        return chart_data

    def _create_chart_object(self) -> None:
        """Creates the chart shape on the slide and populates with data."""
        pos_cfg = self.style.get("position", {})
        dim_cfg = self.style.get("dimensions", {})
        chart_type_str = self.style.get("chart_type", "stacked_bar")

        # Get chart type enum
        if chart_type_str.lower() not in CHART_TYPE_MAP:
            raise ValueError(
                f"Unsupported chart type: '{chart_type_str}'. Supported types are: "
                f"{', '.join(CHART_TYPE_MAP.keys())}"
            )
        xl_chart_type = CHART_TYPE_MAP[chart_type_str.lower()]

        chart_data = self._prepare_chart_data()

        left = Inches(float(pos_cfg.get("left", 1.0)))
        top = Inches(float(pos_cfg.get("top", 1.5)))
        width = Inches(float(dim_cfg.get("width", 8.0)))
        height = Inches(float(dim_cfg.get("height", 5.0)))

        # Add chart shape to slide
        self._chart_shape = self.slide.shapes.add_chart(
            xl_chart_type, left, top, width, height, chart_data
        )
        self.chart = self._chart_shape.chart

    def _apply_chart_title(self) -> None:
        """Applies title settings based on the style configuration."""
        if not self.chart:
            return

        title_config = self.style.get("title", {})
        title_text = title_config.get("text")
        is_visible = title_config.get("visible", True)

        self.chart.has_title = bool(is_visible)

        if self.chart.has_title and title_text is not None:
            title_obj = self.chart.chart_title
            title_obj.has_text_frame = True
            tf = title_obj.text_frame
            tf.text = str(title_text)
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            font_config = title_config.get("font")
            if font_config:
                run = (
                    tf.paragraphs[0].runs[0]
                    if tf.paragraphs[0].runs
                    else tf.paragraphs[0].add_run()
                )
                apply_font_settings(run.font, font_config)

    def _apply_category_axis_settings(self) -> None:
        """Applies settings for the category axis."""
        if not self.chart:
            return

        config = self.style.get("category_axis", {})
        axis = self.chart.category_axis

        if "visible" in config:
            axis.visible = bool(config["visible"])

        if axis.visible:
            if "major_gridlines" in config:
                axis.has_major_gridlines = bool(config["major_gridlines"])
            if "minor_gridlines" in config:
                axis.has_minor_gridlines = bool(config["minor_gridlines"])

            if "tick_marks" in config:
                tick_mark_str = str(config["tick_marks"]).lower()
                axis.major_tick_mark = TICK_MARK_MAP.get(
                    tick_mark_str, TICK_MARK_MAP["none"]
                )

            # Tick Labels
            tick_labels_cfg = config.get("tick_labels", {})
            if "visible" in tick_labels_cfg:
                axis.has_tick_labels = bool(tick_labels_cfg["visible"])

            if axis.has_tick_labels:
                font_config = tick_labels_cfg.get("font")
                if font_config:
                    apply_font_settings(axis.tick_labels.font, font_config)

            # Axis Line Visibility
            line_cfg = config.get("line", {})
            if "visible" in line_cfg and not bool(line_cfg["visible"]):
                axis.format.line.fill.background()

    def _apply_value_axis_settings(self) -> None:
        """Applies settings for the value axis."""
        if not self.chart:
            return

        config = self.style.get("value_axis", {})
        axis = self.chart.value_axis

        if "visible" in config:
            axis.visible = bool(config["visible"])

        if axis.visible:
            if "major_gridlines" in config:
                axis.has_major_gridlines = bool(config["major_gridlines"])
            if "minor_gridlines" in config:
                axis.has_minor_gridlines = bool(config["minor_gridlines"])

            if "tick_marks" in config:
                tick_mark_str = str(config["tick_marks"]).lower()
                axis.major_tick_mark = TICK_MARK_MAP.get(
                    tick_mark_str, TICK_MARK_MAP["none"]
                )

            # Axis Scale Limits
            if "max_scale" in config and config["max_scale"] is not None:
                axis.maximum_scale = float(config["max_scale"])
            if "min_scale" in config and config["min_scale"] is not None:
                axis.minimum_scale = float(config["min_scale"])

            # Tick Labels
            tick_labels_cfg = config.get("tick_labels", {})
            if "visible" in tick_labels_cfg:
                axis.has_tick_labels = bool(tick_labels_cfg["visible"])

            if axis.has_tick_labels:
                if "number_format" in config:
                    axis.tick_labels.number_format = str(config["number_format"])
                    axis.tick_labels.number_format_is_linked = (
                        axis.tick_labels.number_format == "General"
                    )
                else:
                    axis.tick_labels.number_format = "General"
                    axis.tick_labels.number_format_is_linked = True

                font_config = tick_labels_cfg.get("font")
                if font_config:
                    apply_font_settings(axis.tick_labels.font, font_config)

            # Axis Line Visibility
            line_cfg = config.get("line", {})
            if "visible" in line_cfg and not bool(line_cfg["visible"]):
                axis.format.line.fill.background()

    def _apply_plot_area_settings(self) -> None:
        """Applies settings for the plot area (border, fill), if available."""
        if not self.chart or not hasattr(self.chart, "plot_area"):
            return

        config = self.style.get("plot_area", {})
        plot_area = self.chart.plot_area

        # Plot Area Border
        border_cfg = config.get("border", {})
        if "visible" in border_cfg and hasattr(plot_area.format, "line"):
            if not bool(border_cfg["visible"]):
                plot_area.format.line.fill.background()

        # Plot Area Fill
        fill_cfg = config.get("fill", {})
        if "type" in fill_cfg and hasattr(plot_area.format, "fill"):
            fill = plot_area.format.fill
            fill_type = str(fill_cfg["type"]).lower()
            if fill_type == "none":
                fill.background()
            elif fill_type == "solid":
                fill.solid()
                if "color" in fill_cfg:
                    fill.fore_color.rgb = RGBColor.from_string(str(fill_cfg["color"]))
            else:
                raise ValueError(f"Unsupported plot area fill type: '{fill_type}'")

    def _apply_legend_settings(self) -> None:
        """Applies legend settings (position, visibility, font)."""
        if not self.chart:
            return

        config = self.style.get("legend", {})

        if "enabled" in config:
            self.chart.has_legend = bool(config["enabled"])

        if self.chart.has_legend:
            legend = self.chart.legend

            if "position" in config:
                position_str = str(config["position"]).lower()
                legend.position = LEGEND_POSITION_MAP.get(
                    position_str, LEGEND_POSITION_MAP["bottom"]
                )

            if "include_in_layout" in config:
                legend.include_in_layout = bool(config["include_in_layout"])

            font_config = config.get("font")
            if font_config:
                apply_font_settings(legend.font, font_config)

    def _apply_data_labels(self) -> None:
        """Applies data label settings (visibility, position, font, number format)."""
        if not self.chart or not self.chart.plots:
            return

        config = self.style.get("data_labels", {})
        plot = self.chart.plots[0]

        if "enabled" in config:
            plot.has_data_labels = bool(config["enabled"])

        if plot.has_data_labels:
            data_labels = plot.data_labels

            if "position" in config:
                position_str = str(config["position"]).lower()
                data_labels.position = LABEL_POSITION_MAP.get(
                    position_str, LABEL_POSITION_MAP["center"]
                )

            if "number_format" in config:
                data_labels.number_format = str(config["number_format"])
                data_labels.number_format_is_linked = (
                    data_labels.number_format == "General"
                )
            else:
                data_labels.number_format = "General"
                data_labels.number_format_is_linked = True

            font_config = config.get("font")
            if font_config:
                apply_font_settings(data_labels.font, font_config)

    def _apply_series_settings(self) -> None:
        """Applies settings per series, primarily colors."""
        if not self.chart or not self.chart.series:
            return

        series_colors_config = self.style.get("colors", {})

        for i, series in enumerate(self.chart.series):
            # Determine color: Use specific color if defined, else default palette
            color_hex = series_colors_config.get(series.name)
            if not color_hex:
                color_index = i % len(DEFAULT_CHART_COLORS)
                color_hex = DEFAULT_CHART_COLORS[color_index]

            # Apply fill color
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor.from_string(str(color_hex))

    def _apply_bar_column_settings(self) -> None:
        """Applies settings specific to bar/column charts (gap width, overlap)."""
        if not self.chart or not self.chart.plots:
            return

        # Check if the chart type conceptually supports these settings
        supported_types = (
            XL_CHART_TYPE.BAR_STACKED,
            XL_CHART_TYPE.COLUMN_STACKED,
            XL_CHART_TYPE.BAR_CLUSTERED,
            XL_CHART_TYPE.COLUMN_CLUSTERED,
        )

        if self.chart.chart_type in supported_types:
            plot = self.chart.plots[0]

            # Apply gap width if present in style config
            if "gap_width" in self.style:
                plot.gap_width = int(self.style["gap_width"])

            # Apply overlap if present in style config
            if "overlap" in self.style:
                plot.overlap = int(self.style["overlap"])

    def add_to_slide(self) -> Chart:
        """
        Executes the full process: creates the chart object and applies all styles.
        
        Returns:
        --------
        pptx.chart.chart.Chart
            The created and styled PowerPoint chart object.
        """
        self._create_chart_object()

        if not self.chart:
            raise RuntimeError("Chart object was not created successfully.")

        # Apply style settings - order can matter
        self._apply_chart_title()
        self._apply_legend_settings()
        self._apply_category_axis_settings()
        self._apply_value_axis_settings()
        self._apply_plot_area_settings()
        self._apply_series_settings()
        self._apply_data_labels()
        self._apply_bar_column_settings()

        return self.chart


def add_chart(
    slide: Slide, dataframe: pd.DataFrame, style_config: Optional[Dict[str, Any]] = None
) -> Chart:
    """
    High-level function to add a pandas DataFrame as a styled chart to a PowerPoint slide.
    
    Parameters:
    -----------
    slide : pptx.slide.Slide
        The PowerPoint slide object.
    dataframe : pd.DataFrame
        The data to display. Assumes first column is categories, rest are series.
        NaN values will be passed as None, potentially creating gaps in charts.
    style_config : dict, optional
        A dictionary with custom style settings. These will be merged with
        the default style settings.
        
    Returns:
    --------
    pptx.chart.chart.Chart
        The created and styled PowerPoint chart object.
    """
    # Create the final style dictionary
    final_style = create_chart_style(custom_style=style_config)
    
    # Instantiate the formatter and run the process
    formatter = PowerPointChartFormatter(slide, dataframe, final_style)
    
    # Run the creation and styling process
    return formatter.add_to_slide()
