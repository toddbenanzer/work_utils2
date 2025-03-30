"""
Utilities for creating and formatting tables in PowerPoint slides.
"""

from typing import Any, Dict, NamedTuple, Optional, Tuple, Union

import pandas as pd
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.table import Table, _Cell
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> Tuple[int, int, int]:
    """
    Convert hex color string to RGB tuple.

    Args:
        hex_color: Hex color string (e.g., '#4472C4')

    Returns:
        Tuple of (R, G, B) values
    """
    # Remove '#' if present
    hex_color = hex_color.lstrip("#")

    # Convert hex to RGB
    return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))


class TableParams(NamedTuple):
    """Structure to hold calculated table parameters."""

    total_rows: int
    total_cols: int
    index_col_count: int
    data_col_count: int
    header_rows_needed: int
    data_rows_needed: int
    is_multi_index: bool
    is_multi_column: bool


def _calculate_table_parameters(data: pd.DataFrame) -> TableParams:
    """Calculates necessary parameters based on the DataFrame structure."""
    is_multi_index = isinstance(data.index, pd.MultiIndex)
    is_multi_column = isinstance(data.columns, pd.MultiIndex)

    index_col_count = len(data.index.names) if is_multi_index else 1
    data_col_count = len(data.columns)
    total_cols = index_col_count + data_col_count

    header_rows_needed = data.columns.nlevels if is_multi_column else 1
    data_rows_needed = len(data)
    total_rows = header_rows_needed + data_rows_needed

    return TableParams(
        total_rows=total_rows,
        total_cols=total_cols,
        index_col_count=index_col_count,
        data_col_count=data_col_count,
        header_rows_needed=header_rows_needed,
        data_rows_needed=data_rows_needed,
        is_multi_index=is_multi_index,
        is_multi_column=is_multi_column,
    )


def _determine_table_dimensions(
    width: Optional[Union[float, Inches]],
    height: Optional[Union[float, Inches]],
    style_settings: Dict[str, Any],
) -> Tuple[Inches, Inches]:
    """Determines the final width and height for the table shape."""
    actual_width = width
    actual_height = height

    if actual_width is None:
        default_width_inches = style_settings.get("width", 8.0)
        actual_width = Inches(default_width_inches)
        if actual_height is None:
            actual_height = Inches(0.1)  # Auto-height trick
    elif actual_height is None:
        actual_height = Inches(0.1)  # Auto-height trick

    # Ensure Inches type
    if isinstance(actual_width, (int, float)):
        actual_width = Inches(actual_width)
    if isinstance(actual_height, (int, float)):
        actual_height = Inches(actual_height)

    return actual_width, actual_height


def _apply_column_widths(
    table: Table, total_cols: int, style_settings: Dict[str, Any]
) -> None:
    """Applies column widths from style_settings if provided."""
    col_widths = style_settings.get("column_widths")
    if col_widths:
        for i, width_val in enumerate(col_widths):
            if i < total_cols:
                try:
                    # Ensure width_val is treated as inches if it's numeric
                    width_inches = (
                        Inches(width_val)
                        if isinstance(width_val, (int, float))
                        else width_val
                    )
                    table.columns[i].width = width_inches
                except Exception as e:
                    print(f"Warning: Could not set width for column {i}: {e}")


def _fill_headers(table: Table, data: pd.DataFrame, params: TableParams) -> None:
    """Fills the header rows of the table with index and column names."""
    # Fill Index Names (Row 0)
    if params.is_multi_index:
        for i, name in enumerate(data.index.names):
            if name is not None and i < params.index_col_count:
                table.cell(0, i).text = str(name)
    elif data.index.name is not None:
        table.cell(0, 0).text = str(data.index.name)

    # Fill Column Names
    for col_idx, col_value in enumerate(data.columns):
        target_col = col_idx + params.index_col_count
        if target_col >= params.total_cols:
            continue  # Safety break

        if params.is_multi_column:
            # col_value is a tuple (level0_val, level1_val, ...)
            for level, value in enumerate(col_value):
                if level < params.header_rows_needed:
                    table.cell(level, target_col).text = str(value)
        else:
            # col_value is a single name, goes in the first header row
            if 0 < params.header_rows_needed:  # Ensure header row exists
                table.cell(0, target_col).text = str(col_value)


def _fill_data(table: Table, data: pd.DataFrame, params: TableParams) -> None:
    """Fills the data cells of the table."""
    first_data_row_index = params.header_rows_needed

    if params.data_rows_needed <= 0:
        return  # No data rows to fill

    for row_offset, (index_val, data_row) in enumerate(data.iterrows()):
        current_row_index = first_data_row_index + row_offset
        if current_row_index >= params.total_rows:
            break  # Safety break

        # --- Fill index columns for the current data row ---
        if params.is_multi_index:
            if isinstance(index_val, tuple):
                for level, idx_val in enumerate(index_val):
                    if level < params.index_col_count:
                        table.cell(current_row_index, level).text = str(idx_val)
            else:  # Defensive check
                if 0 < params.index_col_count:
                    table.cell(current_row_index, 0).text = str(index_val)
        else:
            table.cell(current_row_index, 0).text = str(index_val)

        # --- Fill data columns for the current data row ---
        col_offset = 0
        # Need to iterate through the *Series* `data_row` correctly
        for value in data_row:  # Iterating directly over the Series gives values
            current_col_index = params.index_col_count + col_offset
            if current_col_index < params.total_cols:
                cell: _Cell = table.cell(current_row_index, current_col_index)
                cell.text = str(value) if pd.notna(value) else ""
            col_offset += 1


def create_table(
    slide,
    data: pd.DataFrame,
    left: Union[float, Inches],
    top: Union[float, Inches],
    width: Optional[Union[float, Inches]] = None,
    height: Optional[Union[float, Inches]] = None,
    style_settings: Optional[Dict[str, Any]] = None,
) -> Optional[Table]:
    """
    Create a table on a slide from DataFrame data.

    Handles single and multi-level indexes and columns with proper formatting.

    Args:
        slide: PowerPoint slide object to add the table to
        data: DataFrame containing the table data
        left: Left position of the table
        top: Top position of the table
        width: Optional width of the table. If None, uses default or auto-calculates
        height: Optional height of the table. If None, uses auto-height
        style_settings: Optional dictionary with style settings

    Returns:
        PowerPoint Table object or None if data is empty
    """
    if style_settings is None:
        style_settings = {}

    # Ensure we're working with a DataFrame
    if not isinstance(data, pd.DataFrame):
        try:
            data = pd.DataFrame(data)
        except Exception as e:
            print(f"Error converting data to DataFrame: {e}")
            return None

    # Handle empty DataFrame
    if data.empty:
        print("Warning: DataFrame is empty, no table created.")
        return None

    # Calculate number of rows and columns needed
    is_multi_index = isinstance(data.index, pd.MultiIndex)
    is_multi_column = isinstance(data.columns, pd.MultiIndex)

    # Determine index columns
    index_col_count = len(data.index.names) if is_multi_index else 1
    data_col_count = len(data.columns)
    total_cols = index_col_count + data_col_count

    # Determine row counts
    header_rows = data.columns.nlevels if is_multi_column else 1
    data_rows = len(data)
    total_rows = header_rows + data_rows

    # Ensure we have some data
    if total_rows < 1 or total_cols < 1:
        print("Warning: Data has no rows or columns.")
        return None

    # Determine table dimensions
    if width is None:
        # Default table width or auto-width (let PowerPoint decide)
        default_width = style_settings.get("width", 8.0)
        table_width = Inches(default_width)
    else:
        # Use specified width
        table_width = width if isinstance(width, Inches) else Inches(width)

    if height is None:
        # Use a very small height to trigger PowerPoint's auto-height behavior
        table_height = Inches(0.1)
    else:
        # Use specified height
        table_height = height if isinstance(height, Inches) else Inches(height)

    # Create the table shape
    try:
        table_shape = slide.shapes.add_table(
            total_rows, total_cols, left, top, table_width, table_height
        )
        table = table_shape.table
    except Exception as e:
        print(f"Error creating table: {e}")
        return None

    # Set column widths if provided
    column_widths = style_settings.get("column_widths")
    if column_widths and isinstance(column_widths, list):
        for i, col_width in enumerate(column_widths):
            if i < total_cols and col_width is not None:
                try:
                    width_val = (
                        Inches(col_width)
                        if isinstance(col_width, (int, float))
                        else col_width
                    )
                    table.columns[i].width = width_val
                except Exception as e:
                    print(f"Warning: Could not set width for column {i}: {e}")

    # Fill header cells (column names)
    if header_rows > 0:
        # Fill index names in header row(s)
        if is_multi_index:
            for i, name in enumerate(data.index.names):
                if name is not None:  # Skip None names (unnamed levels)
                    cell = table.cell(0, i)
                    cell.text = str(name)
        elif data.index.name is not None:  # Single index with a name
            cell = table.cell(0, 0)
            cell.text = str(data.index.name)

        # Fill column headers
        col_offset = index_col_count
        if is_multi_column:
            # Handle multi-level columns
            for col_idx, col_tuple in enumerate(data.columns):
                for level, col_value in enumerate(col_tuple):
                    if level < header_rows:
                        cell = table.cell(level, col_idx + col_offset)
                        cell.text = str(col_value)
        else:
            # Handle single-level columns
            for col_idx, col_name in enumerate(data.columns):
                cell = table.cell(0, col_idx + col_offset)
                cell.text = str(col_name)

    # Fill data cells
    row_offset = header_rows
    for idx, (row_idx, row_data) in enumerate(data.iterrows()):
        # Fill index values
        if is_multi_index:
            # Handle multi-level index
            for level, idx_value in enumerate(row_idx):
                if level < index_col_count:
                    cell = table.cell(idx + row_offset, level)
                    cell.text = str(idx_value)
        else:
            # Handle single-level index
            cell = table.cell(idx + row_offset, 0)
            cell.text = str(row_idx)

        # Fill data values
        for col_idx, value in enumerate(row_data):
            cell = table.cell(idx + row_offset, col_idx + index_col_count)
            # Handle different data types
            if pd.isna(value):
                cell.text = ""
            elif isinstance(value, (int, float)):
                # Format numbers according to settings
                format_str = style_settings.get(f"format_{col_idx}", None)
                if format_str:
                    cell.text = f"{value:{format_str}}"
                else:
                    cell.text = str(value)
            else:
                cell.text = str(value)

    # Apply table styling
    format_table(table, data, style_settings)

    return table


def format_table(
    table: Table, data: pd.DataFrame, style_settings: Optional[Dict[str, Any]] = None
) -> None:
    """
    Apply formatting to a PowerPoint table based on style settings.

    Args:
        table: PowerPoint Table object to format
        data: Original DataFrame (for structure information)
        style_settings: Dictionary with style settings
    """
    if not style_settings:
        return

    # Get basic structure info
    is_multi_index = isinstance(data.index, pd.MultiIndex)
    is_multi_column = isinstance(data.columns, pd.MultiIndex)
    index_col_count = len(data.index.names) if is_multi_index else 1
    header_rows = data.columns.nlevels if is_multi_column else 1

    # Get style settings
    header_style = style_settings.get("header", {})
    default_style = style_settings.get("defaults", {})

    # Format header rows
    if header_rows > 0:
        for row_idx in range(header_rows):
            for col_idx in range(len(table.columns)):
                cell = table.cell(row_idx, col_idx)

                # Apply header formatting
                _apply_cell_format(cell, header_style)

    # Format index columns
    for row_idx in range(header_rows, len(table.rows)):
        for col_idx in range(index_col_count):
            cell = table.cell(row_idx, col_idx)

            # Make index cells slightly different
            index_style = style_settings.get("index", header_style.copy())
            if "fill" in index_style and isinstance(index_style["fill"], str):
                # Use a lighter shade for index
                rgb = hex_to_rgb(index_style["fill"])
                lighter_rgb = tuple(min(int(c * 1.3), 255) for c in rgb)
                index_style["fill"] = (
                    f"#{lighter_rgb[0]:02x}{lighter_rgb[1]:02x}{lighter_rgb[2]:02x}"
                )

            _apply_cell_format(cell, index_style)

    # Format data cells
    for row_idx in range(header_rows, len(table.rows)):
        for col_idx in range(index_col_count, len(table.columns)):
            cell = table.cell(row_idx, col_idx)

            # Check for alternate row styling
            if (
                style_settings.get("banded_rows", False)
                and (row_idx - header_rows) % 2 == 1
            ):
                alt_style = style_settings.get("alternate_row", default_style.copy())
                _apply_cell_format(cell, alt_style)
            else:
                _apply_cell_format(cell, default_style)

            # Apply specific column formatting if available
            col_data_idx = col_idx - index_col_count
            if col_data_idx < len(data.columns):
                col_name = data.columns[col_data_idx]
                col_style = style_settings.get(f"column_{col_data_idx}", {})
                if col_style:
                    _apply_cell_format(cell, col_style)


def _apply_cell_format(cell: _Cell, style: Dict[str, Any]) -> None:
    """
    Apply formatting to a single cell based on style dictionary.

    Args:
        cell: PowerPoint Table Cell object
        style: Dictionary with style settings
    """
    # Text frame formatting
    text_frame = cell.text_frame

    # Set vertical alignment
    v_align = style.get("vertical_alignment")
    if v_align:
        if v_align.lower() == "top":
            cell.vertical_anchor = MSO_ANCHOR.TOP
        elif v_align.lower() == "middle":
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        elif v_align.lower() == "bottom":
            cell.vertical_anchor = MSO_ANCHOR.BOTTOM

    # Process each paragraph in the cell
    for paragraph in text_frame.paragraphs:
        # Set alignment
        alignment = style.get("alignment")
        if alignment:
            if alignment.lower() == "left":
                paragraph.alignment = PP_ALIGN.LEFT
            elif alignment.lower() == "center":
                paragraph.alignment = PP_ALIGN.CENTER
            elif alignment.lower() == "right":
                paragraph.alignment = PP_ALIGN.RIGHT
            elif alignment.lower() == "justify":
                paragraph.alignment = PP_ALIGN.JUSTIFY

        # Font formatting
        font_settings = style.get("font", {})

        # Apply font formatting to all text in paragraph
        for run in paragraph.runs:
            # Font name
            if "name" in font_settings:
                run.font.name = font_settings["name"]

            # Font size
            if "size" in font_settings:
                size_pt = font_settings["size"]
                run.font.size = Pt(size_pt)

            # Bold
            if "bold" in font_settings:
                run.font.bold = font_settings["bold"]

            # Italic
            if "italic" in font_settings:
                run.font.italic = font_settings["italic"]

            # Font color
            if "color" in font_settings:
                color_str = font_settings["color"]
                if color_str.startswith("#"):
                    rgb = hex_to_rgb(color_str)
                    run.font.color.rgb = RGBColor(*rgb)

    # Cell fill (background color)
    fill_color = style.get("fill")
    if fill_color and isinstance(fill_color, str) and fill_color.startswith("#"):
        rgb = hex_to_rgb(fill_color)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*rgb)

    # Cell borders
    border_settings = style.get("border", {})
    if border_settings:
        # Top border
        if border_settings.get("top", False):
            cell.borders.top.width = border_settings.get(
                "width", 12700
            )  # default width

            # Border color
            if "color" in border_settings and border_settings["color"].startswith("#"):
                rgb = hex_to_rgb(border_settings["color"])
                cell.borders.top.color.rgb = RGBColor(*rgb)

        # Bottom border
        if border_settings.get("bottom", False):
            cell.borders.bottom.width = border_settings.get("width", 12700)

            # Border color
            if "color" in border_settings and border_settings["color"].startswith("#"):
                rgb = hex_to_rgb(border_settings["color"])
                cell.borders.bottom.color.rgb = RGBColor(*rgb)

        # Left border
        if border_settings.get("left", False):
            cell.borders.left.width = border_settings.get("width", 12700)

            # Border color
            if "color" in border_settings and border_settings["color"].startswith("#"):
                rgb = hex_to_rgb(border_settings["color"])
                cell.borders.left.color.rgb = RGBColor(*rgb)

        # Right border
        if border_settings.get("right", False):
            cell.borders.right.width = border_settings.get("width", 12700)

            # Border color
            if "color" in border_settings and border_settings["color"].startswith("#"):
                rgb = hex_to_rgb(border_settings["color"])
                cell.borders.right.color.rgb = RGBColor(*rgb)
