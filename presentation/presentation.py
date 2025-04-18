"""
Main module for creating and managing PowerPoint presentations.
Provides classes for creating presentations and slides with various elements.
"""

import json
from typing import Any, Dict, Optional, Union

from pptx import Presentation as PptxPresentation

from . import charts, tables, textbox
from .utils import Pt


class PowerPointPresentation:
    """
    Main class for creating and managing PowerPoint presentations.
    Provides methods for adding slides and saving the presentation.
    """
    
    def __init__(self, template_filepath: Optional[str] = None, config: str = "presentation/config.json"):
        """
        Initialize a PowerPoint presentation.
        
        Parameters:
        -----------
        template_filepath : str, optional
            Path to a template file to use as a base for the presentation.
        config : str, default="presentation/config.json"
            Path to a JSON configuration file with presentation settings.
        """
        self.presentation = PptxPresentation(template_filepath)
        self.config = self._read_config(config)

    def _read_config(self, config_path: str) -> Dict[str, Any]:
        """
        Read configuration from a JSON file.
        
        Parameters:
        -----------
        config_path : str
            Path to the JSON configuration file.
            
        Returns:
        --------
        dict
            The configuration data.
        """
        with open(config_path, "r") as f:
            config_data = json.load(f)
        return config_data

    def add_slide(self, template_index_or_name: Optional[Union[int, str]] = None, config: Optional[Dict[str, Any]] = None) -> 'PowerPointSlide':
        """
        Add a new slide to the presentation.
        
        Parameters:
        -----------
        template_index_or_name : int or str, optional
            Index or name of the slide layout to use.
        config : dict, optional
            Configuration to use for the slide. If None, uses the presentation's config.
            
        Returns:
        --------
        PowerPointSlide
            The newly created slide object.
        """
        if config is None:
            config = self.config
        return PowerPointSlide(self.presentation, template_index_or_name, config)

    def save(self, filepath: str) -> None:
        """
        Save the presentation to a file.
        
        Parameters:
        -----------
        filepath : str
            Path where the presentation will be saved.
        """
        self.presentation.save(filepath)


class PowerPointSlide:
    """
    Class representing a slide in a PowerPoint presentation.
    Provides methods for adding various elements to the slide.
    """
    
    def __init__(self, presentation: PptxPresentation, template_index_or_name: Optional[Union[int, str]] = None, config: Optional[Dict[str, Any]] = None):
        """
        Initialize a PowerPoint slide.
        
        Parameters:
        -----------
        presentation : pptx.Presentation
            The PowerPoint presentation object.
        template_index_or_name : int or str, optional
            Index or name of the slide layout to use.
        config : dict, optional
            Configuration to use for the slide.
        """
        self.presentation = presentation
        self.config = config
        self.slide = self._initialize_slide(template_index_or_name)

    def _get_layout_index_by_name(self, name: str) -> int:
        """
        Get the index of a slide layout by its name.
        
        Parameters:
        -----------
        name : str
            Name of the slide layout.
            
        Returns:
        --------
        int
            Index of the slide layout.
            
        Raises:
        -------
        ValueError
            If the layout name is not found.
        """
        template_name_dict = {
            sl.name: i for i, sl in enumerate(self.presentation.slide_layouts)
        }
        if name in template_name_dict:
            return template_name_dict[name]
        else:
            raise ValueError(
                f"Layout name '{name}' not found. Available layouts: {list(template_name_dict.keys())}"
            )

    def _initialize_slide(self, template_index_or_name: Optional[Union[int, str]] = None) -> Any:
        """
        Initialize a slide with the specified layout.
        
        Parameters:
        -----------
        template_index_or_name : int or str, optional
            Index or name of the slide layout to use.
            
        Returns:
        --------
        pptx.slide.Slide
            The newly created slide object.
            
        Raises:
        -------
        ValueError
            If template_index_or_name is not an int or str.
        """
        if template_index_or_name is None:
            template_index = 0
        elif isinstance(template_index_or_name, int):
            template_index = template_index_or_name
        elif isinstance(template_index_or_name, str):
            template_index = self._get_layout_index_by_name(template_index_or_name)
        else:
            raise ValueError(
                f"template_index_or_name must be an int or str, got {type(template_index_or_name)}"
            )
        slide_layout = self.presentation.slide_layouts[template_index]
        return self.presentation.slides.add_slide(slide_layout)

    def set_title(self, title: str, subtitle: Optional[str] = None) -> None:
        """
        Set the title and optional subtitle of the slide.
        
        Parameters:
        -----------
        title : str
            The title text.
        subtitle : str, optional
            The subtitle text.
        """
        # Set title
        title_placeholder = self.slide.shapes.title
        title_text_frame = title_placeholder.text_frame
        title_p = title_text_frame.paragraphs[0]
        title_p.text = title
        title_p.font.name = self.config["elements"]["title"]["font"]["name"]
        title_p.font.size = Pt(self.config["elements"]["title"]["font"]["size"])
        title_p.font.bold = self.config["elements"]["title"]["font"]["bold"]
        title_p.font.italic = self.config["elements"]["title"]["font"]["italic"]

        # Set subtitle
        if subtitle is not None:
            subtitle_run = title_p.add_run()
            subtitle_run.text = "\n" + subtitle
            subtitle_run.font.name = self.config["elements"]["subtitle"]["font"]["name"]
            subtitle_run.font.size = Pt(
                self.config["elements"]["subtitle"]["font"]["size"]
            )
            subtitle_run.font.bold = self.config["elements"]["subtitle"]["font"]["bold"]
            subtitle_run.font.italic = self.config["elements"]["subtitle"]["font"][
                "italic"
            ]

    def add_textbox(self, text: str, left: float, top: float, width: float, height: float) -> Any:
        """
        Add a textbox to the slide.
        
        Parameters:
        -----------
        text : str
            The text to display in the textbox.
        left : float
            Left position in points.
        top : float
            Top position in points.
        width : float
            Width in points.
        height : float
            Height in points.
            
        Returns:
        --------
        pptx.shapes.autoshape.Shape
            The created textbox shape.
        """
        # Use the textbox module's add_textbox function with custom styling
        style_config = {
            "position": {"left": left / 72, "top": top / 72},  # Convert points to inches
            "dimensions": {"width": width / 72, "height": height / 72},
            "font": {
                "name": self.config["elements"]["textbox"]["font"]["name"],
                "size": self.config["elements"]["textbox"]["font"]["size"],
                "bold": self.config["elements"]["textbox"]["font"]["bold"],
                "italic": self.config["elements"]["textbox"]["font"]["italic"],
            },
        }
        
        return textbox.add_textbox(self.slide, text, style_config)

    def add_table(self, data: Any, style_settings: Optional[Dict[str, Any]] = None) -> Any:
        """
        Add a table to the slide based on DataFrame data.
        
        Parameters:
        -----------
        data : pandas.DataFrame
            DataFrame containing the table data.
        style_settings : dict, optional
            Dictionary with table style settings.
            If None, uses settings from the presentation config file.
            
        Returns:
        --------
        pptx.table.Table
            The created table object.
        """
        return tables.add_table(self.slide, data, style_settings)

    def add_chart(self, data: Any, style_settings: Optional[Dict[str, Any]] = None) -> Any:
        """
        Add a chart to the slide based on DataFrame data.
        
        Parameters:
        -----------
        data : pandas.DataFrame
            DataFrame containing the chart data.
            First column is used as categories (x-axis),
            remaining columns are data series.
        style_settings : dict, optional
            Dictionary with chart style settings.
            If None, uses default chart settings.
            
        Returns:
        --------
        pptx.chart.chart.Chart
            The created chart object.
        """
        return charts.add_chart(self.slide, data, style_settings)
