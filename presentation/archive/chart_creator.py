"""
Utilities for creating and formatting charts in PowerPoint slides.
"""

from typing import Any, Dict, List, Optional, Tuple, Union

import pandas as pd
from pptx.chart.data import CategoryChartData, ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.slide import Slide as PptxSlide
from pptx.util import Inches, Pt

# Mapping of chart types to python-pptx chart types
CHART_TYPE_MAPPING = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "clustered_column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "stacked_column": XL_CHART_TYPE.COLUMN_STACKED,
    "stacked_column_100": XL_CHART_TYPE.COLUMN_STACKED_100,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "clustered_bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "stacked_bar": XL_CHART_TYPE.BAR_STACKED,
    "stacked_bar_100": XL_CHART_TYPE.BAR_STACKED_100,
    "line": XL_CHART_TYPE.LINE,
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "line_stacked": XL_CHART_TYPE.LINE_STACKED,
    "line_stacked_markers": XL_CHART_TYPE.LINE_STACKED_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "area": XL_CHART_TYPE.AREA,
    "stacked_area": XL_CHART_TYPE.AREA_STACKED,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "scatter_smooth": XL_CHART_TYPE.XY_SCATTER_SMOOTH,
}

# Legend position mapping
LEGEND_POSITION_MAPPING = {
    "top": XL_LEGEND_POSITION.TOP,
    "bottom": XL_LEGEND_POSITION.BOTTOM,
    "left": XL_LEGEND_POSITION.LEFT,
    "right": XL_LEGEND_POSITION.RIGHT,
    "corner": XL_LEGEND_POSITION.CORNER,
}


def hex_to_rgb(hex_color: str) -> Tuple[int, int, int]:
    """
    Convert hex color string to RGB tuple.

    Args:
        hex_color: Hex color string (e.g., '#4472C4')

    Returns:
        Tuple of (R, G, B) values
    """
    # Remove '#' if present
    hex_color = hex_color.lstrip("#")

    # Convert hex to RGB
    return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))


def get_chart_data(
    data: pd.DataFrame,
    chart_type: str,
    categories_column: Optional[str] = None,
    values_columns: Optional[List[str]] = None,
) -> ChartData:
    """
    Prepare chart data from DataFrame.

    Args:
        data: DataFrame with chart data
        chart_type: Type of chart
        categories_column: Column to use for categories (x-axis), or None for index
        values_columns: Columns to use for values (y-axis), or None for all numeric

    Returns:
        ChartData object for use with python-pptx
    """
    # Create chart data object
    chart_data = CategoryChartData()

    # Handle categories (x-axis)
    if categories_column is not None and categories_column in data.columns:
        # Use the specified column for categories
        categories = data[categories_column].tolist()
        # Remove the column from consideration for values
        data_for_values = data.drop(columns=[categories_column])
    else:
        # Use index as categories
        categories = data.index.tolist()
        data_for_values = data

    # Set categories
    chart_data.categories = [str(cat) for cat in categories]

    # Handle values (y-axis)
    if values_columns is not None:
        # Use the specified columns
        value_cols = [col for col in values_columns if col in data_for_values.columns]
    else:
        # Use all numeric columns
        numeric_dtypes = ["int64", "float64", "int32", "float32"]
        value_cols = [
            col
            for col in data_for_values.columns
            if data_for_values[col].dtype.name in numeric_dtypes
        ]

    # Add each series
    for col in value_cols:
        values = data_for_values[col].tolist()
        # Convert to float where possible for better chart rendering
        values = [float(v) if pd.notna(v) else 0 for v in values]
        chart_data.add_series(str(col), values)

    return chart_data


def create_chart(
    slide: PptxSlide,
    chart_data: ChartData,
    chart_type: str,
    left: Union[float, Inches],
    top: Union[float, Inches],
    width: Union[float, Inches],
    height: Union[float, Inches],
    style_settings: Optional[Dict[str, Any]] = None,
    chart_title: Optional[str] = None,
):
    """
    Create a chart on a slide.

    Args:
        slide: PowerPoint slide to add chart to
        chart_data: ChartData object with chart data
        chart_type: Type of chart to create
        left: Left position of chart
        top: Top position of chart
        width: Width of chart
        height: Height of chart
        style_settings: Chart style settings from config
        chart_title: Title for the chart or None

    Returns:
        PowerPoint Chart object
    """
    # Get the PowerPoint chart type
    if chart_type in CHART_TYPE_MAPPING:
        pptx_chart_type = CHART_TYPE_MAPPING[chart_type]
    else:
        # Default to clustered column if type not recognized
        pptx_chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED

    # Create the chart
    chart = slide.shapes.add_chart(
        pptx_chart_type, left, top, width, height, chart_data
    ).chart

    # Apply styling if settings provided
    if style_settings:
        _apply_chart_styling(chart, chart_type, style_settings, chart_title)

    return chart


def _apply_chart_styling(
    chart,
    chart_type: str,
    style_settings: Dict[str, Any],
    chart_title: Optional[str] = None,
):
    """
    Apply styling to a chart.

    Args:
        chart: PowerPoint Chart object
        chart_type: Type of chart
        style_settings: Chart style settings from config
        chart_title: Title for the chart or None
    """
    # Set chart title if provided
    if chart_title and style_settings.get("title", {}).get("has_title", True):
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_title

        # Apply title formatting
        title_settings = style_settings.get("title", {})
        if "font_name" in title_settings:
            chart.chart_title.text_frame.paragraphs[0].font.name = title_settings[
                "font_name"
            ]
        if "font_size" in title_settings:
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(
                title_settings["font_size"]
            )
        if "font_bold" in title_settings:
            chart.chart_title.text_frame.paragraphs[0].font.bold = title_settings[
                "font_bold"
            ]
    elif not chart_title:
        chart.has_title = False

    # Set legend properties
    legend_settings = style_settings.get("legend", {})
    chart.has_legend = legend_settings.get("has_legend", True)

    if chart.has_legend:
        # Set legend position
        position = legend_settings.get("position", "bottom").lower()
        if position in LEGEND_POSITION_MAPPING:
            chart.legend.position = LEGEND_POSITION_MAPPING[position]

        # Set legend font
        if "font_name" in legend_settings:
            for paragraph in chart.legend.font.paragraphs:
                paragraph.font.name = legend_settings["font_name"]
        if "font_size" in legend_settings:
            for paragraph in chart.legend.font.paragraphs:
                paragraph.font.size = Pt(legend_settings["font_size"])
        if "font_bold" in legend_settings:
            for paragraph in chart.legend.font.paragraphs:
                paragraph.font.bold = legend_settings["font_bold"]

    # Set axis properties if the chart has axes
    if hasattr(chart, "category_axis") and chart.category_axis:
        axis_settings = style_settings.get("axis", {})
        has_labels = axis_settings.get("has_axis_labels", True)

        # X-axis (category axis)
        chart.category_axis.has_major_gridlines = False
        chart.category_axis.has_minor_gridlines = False
        chart.category_axis.tick_labels.font.name = axis_settings.get(
            "font_name", "Calibri"
        )
        chart.category_axis.tick_labels.font.size = Pt(
            axis_settings.get("font_size", 10)
        )
        chart.category_axis.tick_labels.font.bold = axis_settings.get(
            "font_bold", False
        )

        # Y-axis (value axis)
        if hasattr(chart, "value_axis") and chart.value_axis:
            gridlines = style_settings.get("gridlines", {})
            chart.value_axis.has_major_gridlines = gridlines.get(
                "major_horizontal", True
            )
            chart.value_axis.has_minor_gridlines = gridlines.get(
                "minor_horizontal", False
            )
            chart.value_axis.tick_labels.font.name = axis_settings.get(
                "font_name", "Calibri"
            )
            chart.value_axis.tick_labels.font.size = Pt(
                axis_settings.get("font_size", 10)
            )
            chart.value_axis.tick_labels.font.bold = axis_settings.get(
                "font_bold", False
            )

    # Set data labels
    data_label_settings = style_settings.get("data_labels", {})
    has_data_labels = data_label_settings.get("has_data_labels", False)

    if hasattr(chart, "plots") and chart.plots:
        plot = chart.plots[0]
        plot.has_data_labels = has_data_labels

        if has_data_labels:
            # Format data labels
            data_labels = plot.data_labels
            data_labels.font.name = data_label_settings.get("font_name", "Calibri")
            data_labels.font.size = Pt(data_label_settings.get("font_size", 9))
            data_labels.font.bold = data_label_settings.get("font_bold", False)

            # Set number format for data labels
            number_format = data_label_settings.get("number_format")
            if number_format:
                data_labels.number_format = number_format

    # Set series colors
    colors = style_settings.get("colors", [])
    if colors and hasattr(chart, "series"):
        for i, series in enumerate(chart.series):
            if i < len(colors):
                rgb = hex_to_rgb(colors[i])
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = RGBColor(*rgb)

    # Apply chart type-specific formatting
    if chart_type.startswith("line"):
        _apply_line_chart_styling(chart, style_settings)
    elif chart_type.startswith("column") or chart_type.startswith("bar"):
        _apply_column_bar_styling(chart, style_settings, chart_type)


def _apply_line_chart_styling(chart, style_settings: Dict[str, Any]):
    """
    Apply styling specific to line charts.

    Args:
        chart: PowerPoint Chart object
        style_settings: Chart style settings from config
    """
    # Get line chart specific settings
    line_settings = style_settings.get("line", {})
    line_width = line_settings.get("line_width", 2.5)
    marker_size = line_settings.get("marker_size", 5)
    smooth = line_settings.get("has_smooth_lines", False)

    # Apply to each series
    for series in chart.series:
        if hasattr(series.format, "line"):
            series.format.line.width = Pt(line_width)

        # Set marker size if chart has markers
        if hasattr(series, "marker"):
            series.marker.size = marker_size

        # Set smooth lines
        if hasattr(series, "smooth") and smooth is not None:
            series.smooth = smooth


def _apply_column_bar_styling(chart, style_settings: Dict[str, Any], chart_type: str):
    """
    Apply styling specific to column or bar charts.

    Args:
        chart: PowerPoint Chart object
        style_settings: Chart style settings from config
        chart_type: Type of chart
    """
    # Get column/bar specific settings
    key = "column" if chart_type.startswith("column") else "bar"
    settings = style_settings.get(key, {})

    # Set gap width (space between clusters)
    gap_width = settings.get("gap_width", 150)

    # Set overlap (space between bars in a cluster)
    overlap = settings.get("overlap", 0)

    # Apply to chart
    if hasattr(chart, "plots") and chart.plots:
        plot = chart.plots[0]
        if hasattr(plot, "gap_width"):
            plot.gap_width = gap_width
        if hasattr(plot, "overlap") and chart_type.startswith(
            ("column_clustered", "bar_clustered")
        ):
            plot.overlap = overlap
